# create_presentation.py
#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import pandas as pd
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime

def create_presentation_with_error_stats(error_stats_data, output_dir):
    """
    Создает полноценную PowerPoint презентацию со статистикой ошибок
    """
    # Создаем новую презентацию
    prs = Presentation()
    
    # СЛАЙД 1: Титульный слайд
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Анализ точности прогнозирования\nбереговой эрозии"
    subtitle.text = f"Статистика ошибок модели MLP\n{datetime.now().strftime('%d.%m.%Y')}"
    
    # СЛАЙД 2: Общая информация
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд
    
    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Общая информация"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    # Содержание
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    total_sites = len(error_stats_data)
    total_observations = sum(error_stats_data['Количество_наблюдений'])
    avg_abs_error = error_stats_data['Средняя_абс_ошибка'].mean()
    
    content_text = f"""
    📊 ОБЗОР ДАННЫХ:
    
    • Количество участков: {total_sites}
    • Общее количество наблюдений: {total_observations}
    • Средняя абсолютная ошибка: {avg_abs_error:.2f} м
    • Период анализа: {datetime.now().strftime('%d.%m.%Y')}
    
    🎯 МЕТОДОЛОГИЯ:
    
    • Модель: MLP (многослойный перцептрон)
    • Целевая переменная: Расстояние до ПН
    • Признаки: Год, Участок, Профиль, Литология
    • Горизонт прогнозирования: 20 лет
    
    📈 КРИТЕРИИ ОЦЕНКИ:
    
    • Высокая точность: < 5 м
    • Средняя точность: 5-10 м  
    • Низкая точность: ≥ 10 м
    """
    
    content_frame.text = content_text
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.name = "Arial"
    
    # СЛАЙД 3: Таблица статистики ошибок
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Статистика ошибок по участкам"
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    # Создаем таблицу
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5)
    
    rows = len(error_stats_data) + 1
    cols = len(error_stats_data.columns)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Заголовки столбцов
    headers = ['Участок', 'Средняя\nошибка', 'Стд\nошибки', 'Мин\nошибка', 
               'Макс\nошибка', 'Средняя\nабс. ошибка', 'Наблюдения']
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.bold = True
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 78, 120)
    
    # Заполняем данные
    for row_idx, (_, row) in enumerate(error_stats_data.iterrows(), 1):
        for col_idx, col_name in enumerate(error_stats_data.columns):
            cell = table.cell(row_idx, col_idx)
            
            if col_name in ['Средняя_ошибка', 'Стд_ошибка', 'Мин_ошибка', 
                          'Макс_ошибка', 'Средняя_абс_ошибка']:
                value = f"{row[col_name]:.2f}"
            elif col_name == 'Количество_наблюдений':
                value = f"{int(row[col_name])}"
            else:
                value = str(row[col_name])
            
            cell.text = value
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(10)
            
            # Цветовое кодирование для средней абсолютной ошибки
            if col_name == 'Средняя_абс_ошибка':
                if row[col_name] < 5:
                    paragraph.font.color.rgb = RGBColor(0, 128, 0)  # Зеленый
                    paragraph.font.bold = True
                elif row[col_name] < 10:
                    paragraph.font.color.rgb = RGBColor(255, 165, 0)  # Оранжевый
                else:
                    paragraph.font.color.rgb = RGBColor(220, 0, 0)  # Красный
                    paragraph.font.bold = True
            
            # Чередование цветов фона
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 249, 250)
    
    # СЛАЙД 4: Визуализация точности
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Визуализация точности прогнозирования"
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    # Создаем график
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))
    
    # График 1: Средняя абсолютная ошибка по участкам
    sites = error_stats_data['Участок']
    abs_errors = error_stats_data['Средняя_абс_ошибка']
    
    colors = ['green' if err < 5 else 'orange' if err < 10 else 'red' for err in abs_errors]
    
    bars = ax1.bar(range(len(sites)), abs_errors, color=colors, alpha=0.7)
    ax1.set_title('Средняя абсолютная ошибка по участкам', fontsize=12, fontweight='bold')
    ax1.set_ylabel('Ошибка (м)')
    ax1.set_xticks(range(len(sites)))
    ax1.set_xticklabels(sites, rotation=45, ha='right', fontsize=8)
    ax1.grid(True, alpha=0.3)
    
    # Добавляем значения на столбцы
    for bar, error in zip(bars, abs_errors):
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height + 0.1,
                f'{error:.1f}', ha='center', va='bottom', fontsize=8, fontweight='bold')
    
    # График 2: Распределение по категориям точности
    categories = ['Высокая (<5м)', 'Средняя (5-10м)', 'Низкая (≥10м)']
    counts = [
        len(error_stats_data[error_stats_data['Средняя_абс_ошибка'] < 5]),
        len(error_stats_data[(error_stats_data['Средняя_абс_ошибка'] >= 5) & 
                           (error_stats_data['Средняя_абс_ошибка'] < 10)]),
        len(error_stats_data[error_stats_data['Средняя_абс_ошибка'] >= 10])
    ]
    colors_pie = ['#28a745', '#ffc107', '#dc3545']
    
    ax2.pie(counts, labels=categories, autopct='%1.1f%%', colors=colors_pie,
           startangle=90, textprops={'fontsize': 10})
    ax2.set_title('Распределение участков по точности', fontsize=12, fontweight='bold')
    
    plt.tight_layout()
    
    # Сохраняем график временно
    chart_path = os.path.join(output_dir, "temp_chart.png")
    plt.savefig(chart_path, dpi=150, bbox_inches='tight')
    plt.close()
    
    # Добавляем график в презентацию
    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.3), Inches(9), Inches(4.5))
    
    # Удаляем временный файл
    os.remove(chart_path)
    
    # СЛАЙД 5: Анализ лучших и худших результатов
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Детальный анализ результатов"
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    # Находим лучшие и худшие участки
    best_idx = error_stats_data['Средняя_абс_ошибка'].idxmin()
    worst_idx = error_stats_data['Средняя_абс_ошибка'].idxmax()
    
    best_site = error_stats_data.loc[best_idx]
    worst_site = error_stats_data.loc[worst_idx]
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    analysis_text = f"""
    🏆 ЛУЧШИЕ РЕЗУЛЬТАТЫ:
    
    • Участок: {best_site['Участок']}
    • Средняя абсолютная ошибка: {best_site['Средняя_абс_ошибка']:.2f} м
    • Количество наблюдений: {best_site['Количество_наблюдений']}
    • Стандартное отклонение: {best_site['Стд_ошибка']:.2f} м
    
    ⚠️ НАИБОЛЕЕ ПРОБЛЕМНЫЕ УЧАСТКИ:
    
    • Участок: {worst_site['Участок']}
    • Средняя абсолютная ошибка: {worst_site['Средняя_абс_ошибка']:.2f} м  
    • Количество наблюдений: {worst_site['Количество_наблюдений']}
    • Стандартное отклонение: {worst_site['Стд_ошибка']:.2f} м
    
    📋 РЕКОМЕНДАЦИИ:
    
    • Увеличить сбор данных для участков с низкой точностью
    • Проверить качество данных проблемных участков
    • Рассмотреть дополнительные факторы для сложных участков
    • Использовать участки с высокой точностью для калибровки модели
    """
    
    content_frame.text = analysis_text
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.name = "Arial"
    
    # СЛАЙД 6: Заключение
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Выводы и заключение"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    conclusion_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4))
    conclusion_frame = conclusion_box.text_frame
    conclusion_frame.word_wrap = True
    
    high_accuracy = len(error_stats_data[error_stats_data['Средняя_абс_ошибка'] < 5])
    medium_accuracy = len(error_stats_data[(error_stats_data['Средняя_абс_ошибка'] >= 5) & 
                                         (error_stats_data['Средняя_абс_ошибка'] < 10)])
    low_accuracy = len(error_stats_data[error_stats_data['Средняя_абс_ошибка'] >= 10])
    
    conclusion_text = f"""
    📊 ОСНОВНЫЕ ВЫВОДЫ:
    
    • Модель демонстрирует хорошую точность для {high_accuracy} из {len(error_stats_data)} участков
    • Средняя точность достигнута для {medium_accuracy} участков
    • Требуют внимания {low_accuracy} участков с низкой точностью
    
    🎯 ПЕРСПЕКТИВЫ РАЗВИТИЯ:
    
    • Оптимизация гиперпараметров модели
    • Включение дополнительных признаков
    • Увеличение объема обучающих данных
    • Внедрение ансамблевых методов
    
    💡 ПРАКТИЧЕСКОЕ ПРИМЕНЕНИЕ:
    
    • Результаты могут быть использованы для планирования берегозащитных мероприятий
    • Модель позволяет прогнозировать динамику береговой линии на 20 лет
    • Рекомендуется регулярное обновление модели по мере поступления новых данных
    """
    
    conclusion_frame.text = conclusion_text
    for paragraph in conclusion_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.name = "Arial"
    
    # Сохраняем презентацию
    output_path = os.path.join(output_dir, "error_analysis_presentation.pptx")
    prs.save(output_path)
    
    print(f"✅ Презентация успешно создана: {output_path}")
    print(f"📊 Количество слайдов: {len(prs.slides)}")
    print(f"🏷️ Участков в анализе: {len(error_stats_data)}")

def main():
    """Основная функция"""
    # Данные статистики ошибок
    error_stats_data = {
        'Участок': ['Бережновка', 'Бурты', 'Лист1', 'Молчановка', 'Нижний Балыклей', 
                   'Нижний Ураков', 'Новоникольское', 'Пичуга-Южный', 'Пролейский', 
                   'Суводская', 'Ураков Бугор'],
        'Средняя_ошибка': [1.387, 0.12, 2.197, 0.976, -0.026, 0.643, -1.763, 1.432, 0.429, 0.26, 0.271],
        'Стд_ошибка': [18.786, 3.786, 59.334, 7.006, 9.318, 4.403, 26.441, 5.944, 5.166, 2.435, 6.264],
        'Мин_ошибка': [-38.299, -8.3, -98.34, -18.775, -15.067, -9.347, -37.324, -11.913, -12.003, -7.483, -20.775],
        'Макс_ошибка': [75.748, 6.884, 81.198, 9.826, 21.86, 10.058, 79.221, 13.903, 9.98, 3.272, 13.198],
        'Средняя_абс_ошибка': [14.54, 3.112, 42.686, 5.752, 8.019, 3.47, 21.082, 4.231, 4.467, 1.368, 4.751],
        'Количество_наблюдений': [71, 74, 8, 71, 68, 86, 56, 43, 68, 14, 63]
    }
    
    # Создаем DataFrame
    error_df = pd.DataFrame(error_stats_data)
    
    # Папка для результатов
    output_dir = "presentation_results"
    os.makedirs(output_dir, exist_ok=True)
    
    # Создаем презентацию
    create_presentation_with_error_stats(error_df, output_dir)

if __name__ == "__main__":
    main()
